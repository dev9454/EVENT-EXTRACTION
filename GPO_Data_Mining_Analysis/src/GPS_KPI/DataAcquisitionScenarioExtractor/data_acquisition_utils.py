from pathlib import Path
import pickle
from datetime import datetime
import geopandas as gpd
import numpy as np
import pandas as pd
import pptx
import shapely
from pptx import Presentation
from pptx.util import Pt
from dataclasses import dataclass
import warnings

warnings.simplefilter(action='ignore', category=FutureWarning)


@dataclass
class Area:
    min_x: float
    min_y: float
    max_x: float
    max_y: float


@dataclass
class RoadType:
    name: str
    area: Area
    raw_file_path: str
    processed_file_path: str
    processed_merged_file_path: str


@dataclass
class Config:
    config_path: str
    region: str
    administrative_area: RoadType = None
    barrier: RoadType = None
    bridges: RoadType = None
    buildings: RoadType = None
    bus_stop: RoadType = None
    crossing: RoadType = None
    gantry: RoadType = None
    intersections: RoadType = None
    manhole: RoadType = None
    landuse: RoadType = None
    parking: RoadType = None
    pole: RoadType = None
    power_line: RoadType = None
    railway_crossing: RoadType = None
    road_type: RoadType = None
    roundabout: RoadType = None
    route_relations: RoadType = None
    street_lamp: RoadType = None
    toll_booth: RoadType = None
    traffic_calming: RoadType = None
    traffic_sign: RoadType = None
    traffic_signals: RoadType = None
    tunnel: RoadType = None


def read_file(file_input_path: str) -> gpd.GeoDataFrame:
    """Read a file and return a GeoDataFrame.

       Args:
           file_input_path (str): The path to the input file,
               which must be a .gpkg or .pkl file.

       Returns:
           gpd.GeoDataFrame
       """
    file_input_path = Path(file_input_path)
    if file_input_path.is_file() and file_input_path.suffix == '.gpkg':
        return gpd.read_file(file_input_path)
    elif file_input_path.is_file() and file_input_path.suffix == '.pkl':
        with open(file_input_path, 'rb') as f:
            gdf = pd.read_pickle(f)
        # Convert to GeoDataFrame if needed (only if it's not already a GeoDataFrame)
        if not isinstance(gdf, gpd.GeoDataFrame):
            gdf = gpd.GeoDataFrame(gdf)
        return gdf
    else:
        raise ValueError("File must be a .gpkg or .pkl file.\n Current file format is: " + file_input_path.suffix)


def geopoint_to_dms(point: shapely.geometry.Point) -> str:
    """Convert a geographic point to degrees, minutes, and seconds.

    Args:
        point (shapely.geometry.Point): A point representing geographic coordinates.

    Returns:
        str: Latitude and longitude in DMS format.
        For example:
        '37°47'22.12"N\n122°25'9.37"W'.
    """
    if point is None:
        return None
    coords = shapely.get_coordinates(point)
    long = coords[0, 0]
    lat = coords[0, 1]
    if lat > 0:
        lat_dir = "N"
    else:
        lat_dir = "S"

    if long > 0:
        long_dir = "E"
    else:
        long_dir = "W"

    lat = abs(lat)
    long = abs(long)

    lat_d = int(lat)
    lat_m = int((lat - lat_d) * 60)
    lat_s = (lat - lat_d - lat_m / 60) * 3600

    long_d = int(long)
    long_m = int((long - long_d) * 60)
    long_s = (long - long_d - long_m / 60) * 3600

    return f"{lat_d}°{lat_m}'{lat_s:.2f}\"{lat_dir}\n{long_d}°{long_m}'{long_s:.2f}\"{long_dir}"


def combine_columns(row: pd.Series, columns: list) -> str:
    """Combine columns to create a comment

    Args:
        row (pandas.Series)
        columns (list): List of column names to combine.

    Returns:
        str: A string with all comments combined
    """
    comment = ''
    for column in columns:
        line = f"{column} = {str(row[column]) if row[column] else 'Unknown'}\n"
        comment += line
    return comment.strip()


def create_pptx_df(df: gpd.GeoDataFrame, comments_columns: list) -> gpd.GeoDataFrame:
    """Creates a dataframe that can be used to fill the pptx file

    Args:
        df (gpd.GeoDataFrame)
        comments_columns (list): List of column names to combine.

    Returns:
        (gpd.GeoDataFrame): A dataframe for pptx creation
    """
    col_1 = ['No.', 'Landmark', 'Start', 'Stop', 'Comments']
    df1 = pd.DataFrame(columns=col_1)
    no1 = [i + 1 for i in range(len(df.index))]
    df1['No.'] = no1
    try:
        df1['Landmark'] = df['name']
    except KeyError:
        # Try to get name_1 to handle cases like intersections
        try:
            df1['Landmark'] = df['name_1']
        except:
            df1['Landmark'] = ''
    df1['Comments'] = df.apply(combine_columns, axis=1, args=(comments_columns,))
    df1["Start"], df1["Stop"] = zip(*df.geometry.apply(get_start_stop))
    df1["Start"] = df1["Start"].apply(geopoint_to_dms)
    df1["Stop"] = df1["Stop"].apply(geopoint_to_dms)
    return df1.replace('', np.nan)


from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor


def create_pptx(input_df: gpd.GeoDataFrame, output_file: Path | str, comments_columns: list) -> None:
    """Creates and saves a pptx presentation to a file, with a maximum of 10 slides.

    Args:
        input_df (gpd.GeoDataFrame): DataFrame with content to add to slides.
        output_file (str/Path): Path to output file.
        comments_columns (list): List of column names to combine.
    """
    template_file = r"/mnt/usmidet/users/d0ir4n/template_1.pptx"
    # Create comments column
    pptx_df = create_pptx_df(input_df, comments_columns)
    # Load the presentation template
    prs = Presentation(template_file)
    prs.slides[0].shapes[0].text = "HOV Lanes"
    prs.slides[0].shapes[1].text = datetime.today().strftime('%Y/%m/%d')
    for i in range(min((len(pptx_df.index) // 6) + 1, 10)):
        try:
            slide = prs.slides[i + 1]  # +1 since first slide i
        except IndexError:
            continue
        # Find the table in the slide
        table = None
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                break
        if table is None:
            raise ValueError("No table found in the specified slide.")
        for col_idx, col_name in enumerate(pptx_df.columns):
            cell = table.cell(0, col_idx)
            cell.text = col_name
            cell.text_frame.paragraphs[0].font.bold = True
        # Update table content
        for row_idx, row in pptx_df[i * 6:i * 6 + 6].iterrows():
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx % 6 + 1, col_idx)
                cell.text = str(value)

                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = pptx.dml.color.RGBColor(255, 255, 255)

                if col_idx == 4:
                    for para in cell.text_frame.paragraphs:
                        para.runs[0].font.color.rgb = pptx.dml.color.RGBColor(255, 255, 255)
                        para.runs[0].font.size = Pt(11)
    prs.save(output_file)


def create_map(gdf: gpd.GeoDataFrame, output_file_path: str) -> None:
    """
    Creates an html file containing map with plotted objects

    Args:
    gdf (gpd.GeoDataFrame)
    output_file_path (str): A path to output
    """
    gdf.explore().save(output_file_path)


def save_file(gdf: gpd.GeoDataFrame, output_file_path: str, save_as_excel: bool = True) -> None:
    """
    Saves a GeoDataFrame to a .gpkg file and optionally to an Excel file.

    Args:
        gdf (gpd.GeoDataFrame): The GeoDataFrame to save.
        output_file_path (str): The path to the .gpkg output file.
        save_as_excel (bool): If True, also saves a copy as an Excel file (.xlsx).
    """

    if not str(output_file_path).endswith(".gpkg"):
        output_file_path += ".gpkg"

    gdf.to_file(output_file_path, driver="GPKG")

    if save_as_excel:
        start_points = gdf.geometry.apply(get_start_stop).str[0]
        dms = start_points.apply(geopoint_to_dms).str.replace("\n", "+").str.replace("\\", "")
        google_links = "https://www.google.pl/maps/place/" + dms
        gdf["google_maps_link"] = google_links
        excel_path = str(output_file_path).replace(".gpkg", ".xlsx")

        # Excel limits number of URLS to 65530. Convert URL to strings
        writer = pd.ExcelWriter(
            excel_path,
            engine='xlsxwriter',
            engine_kwargs={'options': {'strings_to_urls': False}}
        )

        gdf.to_excel(writer, index=False)
        writer.close()


def get_intersection(gdf_a: gpd.GeoDataFrame, gdf_b: gpd.GeoDataFrame) -> gpd.GeoDataFrame:
    """Function to calculate intersection points of two geopandas dataframe. Function check if gdf_b intersects with gdf_a

    Args:
        gdf_a (GeoDataFrame): First GeoDataFrame
        gdf_b (GeoDataFrame): Second GeoDataFrame

    Returns:
        GeoDataFrame: GeoDataFrame with intersection points. This datafrme has columns from gdf_a, gdf_b and inter_geometry as active geometry
    """
    gdf_b = gdf_b.copy()
    gdf_b['b_geometry'] = gdf_b['geometry']
    intersects = gpd.sjoin(left_df=gdf_a, right_df=gdf_b, how="inner", predicate="intersects").reset_index(drop=True)

    #  get intersection for each row
    intersects['inter_geometry'] = pd.Series(map(shapely.intersection,
                                                 intersects['geometry'],
                                                 intersects['b_geometry']))

    intersects = intersects.set_geometry('inter_geometry')
    #  explode rows which are multipoints / multilinestrings
    intersects_exp = intersects.explode(ignore_index=True)
    return intersects_exp.set_crs(4326)


def merge_gdf_information(gdf_a: gpd.GeoDataFrame, gdf_b: gpd.GeoDataFrame, gdf_a_name: str = 'left',
                          gdf_b_name: str = 'right', max_distance: float = 1e-4) -> gpd.GeoDataFrame:
    """Function to merge information from both GeoDataFrames. Function search for each object from gdf_a
    the closest/crossing object in gdf_b (if max distance is lower than max_distance value).

    Args:
        gdf_a (GeoDataFrame): First GeoDataFrame
        gdf_b (GeoDataFrame): Second GeoDataFrame
        gdf_a_name (str): First GeoDataFrame name
        gdf_b_name (str): Second GeoDataFrame name
        max_distance (float): Max distance between object from First and Second GeoDataFrames
                              (for crossing objects leave default value)

    Returns:
        GeoDataFrame: First GeoDataFrame merged with the closest/crossing object from Second GeoDataFrame.
        This dataframe has columns from gdf_a and gdf_b.
    """

    gdf_b = gdf_b.copy()
    gdf_b[f'geometry_{gdf_b_name}'] = gdf_b['geometry']
    merged_gdf = gpd.sjoin_nearest(left_df=gdf_a.to_crs(3857), right_df=gdf_b.to_crs(3857), how='left',
                                   lsuffix=gdf_a_name, rsuffix=gdf_b_name, max_distance=max_distance)
    #  Remove duplicated rows
    merged_gdf.drop_duplicates(subset=[f'el_id_{gdf_a_name}'], inplace=True)
    return merged_gdf.to_crs(4326)


def extract_directory(file_path: str) -> Path:
    """Returns directory path of a file"""
    return Path(file_path).parent


def get_start_stop(geometry):
    """
    Gets start and stop based on geometry type
    Args:
        geometry:

    Returns:
        start(shapely.Point): Start point of geometry
        stop(shapely.Point): Stop point of geometry
    """
    if geometry.geom_type == "LineString":
        # Get start and stop points of LineString
        start = shapely.get_point(geometry, 0)
        stop = shapely.get_point(geometry, -1)
    elif geometry.geom_type == "Point":
        # For point, assign point to start and leave stop as empty
        start = geometry
        stop = None
    elif geometry.geom_type in ["Polygon", "MultiPolygon"]:
        # For polygon shapes, get centroid
        start = geometry.centroid
        stop = None
    elif geometry.geom_type == "MultiLineString":
        # For MultiLineString, get start as first point of first LineString
        # and stop as last point of last LineString
        start = shapely.get_point(geometry, 0, 0)
        stop = shapely.get_point(geometry, -1, -1)
    else:
        start = None
        stop = None
    return start, stop